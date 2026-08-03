"""자표준문서.xlsx 의 짝이 되는 **합성 대상 문서 + 골든셋**을 만든다 (Phase F3.5).

`samples/자표준문서.xlsx` 는 배터리 규격 기준 문서인데 대조할 대상 문서가 없어
라이브 검증·매칭 spike·F6 벤치마크를 돌릴 수 없었다. 이 스크립트는 **하나의 CASES
테이블**에서 대상 문서(docx/pptx)와 정답(골든셋 jsonl)을 **동시에** 생성한다 —
문서와 정답이 같은 소스에서 나오므로 정답 드리프트가 원천적으로 불가능하다.

산출물:
    samples/자표준_규격서.docx      Word 대상(문단 + 표)
    samples/자표준_발표.pptx        PPT 대상(텍스트박스 + 표 + 스피커노트)
    golden/자표준_골든셋.jsonl      entity×문서 단위 정답(match/mismatch/missing/unknown)

실행:
    python scripts/make_synthetic_targets.py            # 기본 경로에 생성
    python scripts/make_synthetic_targets.py --out-dir tmp

의존성: python-docx, python-pptx (COM 불필요 — 어느 OS 에서나 생성 가능).
읽기(추출)는 여전히 Windows + Office COM 이 필요하다.
"""

from __future__ import annotations

import argparse
import json
import os
import sys

# --------------------------------------------------------------------------- #
# 기준 문서(samples/자표준문서.xlsx)의 실제 값 — 행/열 좌표까지 그대로 옮긴다.
# 헤더는 2행 구조: 1행 "정량규격", 2행 F/G/H = 하한치/중심치/상한치.
# ⚠ 이 시트는 단위 열(K)이 전부 비어 있고, 단일값 항목이 F(하한치) 또는 G(중심치)에
#   섞여 들어가 있다 — 실문서의 지저분함을 그대로 둔 채 정답을 매긴다.
# --------------------------------------------------------------------------- #
REF: dict[str, dict] = {
    "고객 표준 버전": {"row": 3, "cells": "E3:J3",
                       "attrs": {"qualitative_spec": "배터리승인규격 ver 4.7 SEC Req. ver.4.7"}},
    "공칭용량": {"row": 4, "cells": "E4:F4", "attrs": {"lower_limit": 1150}},
    "정격용량": {"row": 5, "cells": "E5:F5", "attrs": {"lower_limit": 1150}},
    "정격충전전압": {"row": 6, "cells": "E6:G6", "attrs": {"target_value": 4.55}},
    "공칭전압": {"row": 7, "cells": "E7:G7", "attrs": {"target_value": 3.89}},
    "deltaOCV": {"row": 8, "cells": "E8", "attrs": {}},
    "충전방법": {"row": 9, "cells": "E9", "attrs": {}},
    "충전방법(SOC28->SOC64)": {"row": 10, "cells": "E10", "attrs": {}},
    "표준충전전류": {"row": 11, "cells": "E11:G11", "attrs": {"target_value": 230}},
    "최대충전전류": {"row": 12, "cells": "E12:G12", "attrs": {"target_value": 1495}},
    "최대방전전류": {"row": 13, "cells": "E13:G13", "attrs": {"target_value": 1150}},
    "방전종지전압": {"row": 14, "cells": "E14:G14", "attrs": {"target_value": 3}},
    "표준환경온도": {"row": 15, "cells": "E15:H15",
                     "attrs": {"lower_limit": 21, "target_value": 25, "upper_limit": 29}},
    "평가환경습도": {"row": 16, "cells": "E16:H16",
                     "attrs": {"lower_limit": 33, "target_value": 43, "upper_limit": 53}},
    "충전환경온도": {"row": 17, "cells": "E17:H17",
                     "attrs": {"lower_limit": -5, "target_value": 35, "upper_limit": 85}},
    "방전환경온도": {"row": 18, "cells": "E18:H18",
                     "attrs": {"lower_limit": -30, "target_value": 30, "upper_limit": 90}},
    "평가환경온도": {"row": 19, "cells": "E19:H19",
                     "attrs": {"lower_limit": 21, "target_value": 24, "upper_limit": 28}},
    "1개월저장온도": {"row": 20, "cells": "E20:H20",
                      "attrs": {"lower_limit": -10, "target_value": 35, "upper_limit": 80}},
    "3개월저장온도": {"row": 21, "cells": "E21:H21",
                      "attrs": {"lower_limit": -10, "target_value": 35, "upper_limit": 70}},
    "1년저장온도": {"row": 22, "cells": "E22:H22",
                    "attrs": {"lower_limit": -10, "target_value": 35, "upper_limit": 55}},
}

WORD_NAME = "자표준_규격서.docx"
PPT_NAME = "자표준_발표.pptx"

# --------------------------------------------------------------------------- #
# CASES — 대상 문서에 어떻게 쓸지 + 그 결과 정답이 무엇인지를 한 줄에 담는다.
#
#   entity        기준 문서의 항목명
#   text          대상 문서에 실제로 들어갈 문구(표 행이면 셀 목록)
#   expected      match | mismatch | missing | unknown  (계획 §6.2 의 4분류)
#   bad           mismatch 일 때 어긋난 attribute 이름들
#   reason        사람이 그렇게 라벨한 이유(검수·회고용)
#
# 라벨 기준(기준 문서에 단위 열이 비어 있어 필요한 약속):
#   · 값이 같고 대상 단위가 그 항목의 상식적 단위면      → match
#   · 값 자체가 다르면                                    → mismatch
#   · 값은 같아 보이나 **단위 스케일이 달라**(1495 vs 1.495A)
#     기준의 단위를 모르면 등가를 확정할 수 없으면        → unknown
#   · 대상 문서에 항목 자체가 없으면                      → missing
# --------------------------------------------------------------------------- #
WORD_CASES: list[dict] = [
    {"entity": "고객 표준 버전", "where": "para",
     "text": "본 규격은 배터리승인규격 ver 4.7 (SEC Req. ver.4.7) 을 따른다.",
     "expected": "match", "reason": "정성 규격 문구가 기준과 동일"},
    # 한 문장에 두 항목을 담으면 "어느 fact 가 이 항목의 정답인지"가 모호해져
    # 채점이 흔들린다(문서 품질이 아니라 정답 귀속의 문제). 항목당 한 문장으로 둔다.
    {"entity": "정격충전전압", "where": "para",
     "text": "정격 충전 전압은 4.55V 이다.",
     "expected": "match", "reason": "4.55 일치"},
    {"entity": "방전종지전압", "where": "para",
     "text": "방전 종지 전압은 3.0V 로 한다.",
     "expected": "match", "reason": "기준 3 과 표기만 다름(3 vs 3.0)"},
    {"entity": "공칭전압", "where": "para",
     "text": "공칭전압은 3.85V 이다.",
     "expected": "mismatch", "bad": ["target_value"], "reason": "기준 3.89 vs 대상 3.85"},
    {"entity": "최대충전전류", "where": "para",
     "text": "최대 충전 전류는 1.495A 를 초과하지 않는다.",
     "expected": "unknown",
     "reason": "기준 1495 는 단위 열이 비어 있어 mA 인지 확정 불가 — 1.495A 와 등가인지 판단보류"},
    {"entity": "평가환경습도", "where": "para",
     "text": "평가 환경 습도는 33~53%RH 범위로 하며 중심치는 43%RH 이다.",
     "expected": "match", "reason": "33/43/53 모두 일치"},
    # --- 표 ---
    {"entity": "공칭용량", "where": "table", "row": ["공칭용량", "1150", "mAh"],
     "expected": "match", "reason": "1150 일치(기준은 하한치 열에 단일값이 들어가 있음)"},
    {"entity": "정격용량", "where": "table", "row": ["정격용량", "1150", "mAh"],
     "expected": "match", "reason": "1150 일치"},
    {"entity": "표준충전전류", "where": "table", "row": ["표준충전전류", "230", "mA"],
     "expected": "match", "reason": "230 일치"},
    {"entity": "최대방전전류", "where": "table", "row": ["최대방전전류", "1150", "mA"],
     "expected": "match", "reason": "1150 일치"},
    {"entity": "표준환경온도", "where": "table", "row": ["표준환경온도", "21 ~ 29 (중심 25)", "℃"],
     "expected": "match", "reason": "21/25/29 모두 일치"},
    # --- 대상 문서에 아예 없는 항목 ---
    {"entity": "deltaOCV", "where": "absent", "expected": "missing",
     "reason": "규격서에 언급 없음"},
    {"entity": "충전방법", "where": "absent", "expected": "missing",
     "reason": "규격서에 언급 없음"},
    {"entity": "1년저장온도", "where": "absent", "expected": "missing",
     "reason": "저장 조건은 발표자료에만 있음"},
]

PPT_CASES: list[dict] = [
    {"entity": "고객 표준 버전", "where": "s1",
     "text": "적용 규격: SEC Req. ver.4.7 (배터리승인규격 ver 4.7)",
     "expected": "match", "reason": "정성 규격 문구 일치"},
    # slide 2 — 환경 조건 표
    {"entity": "충전환경온도", "where": "s2", "row": ["충전환경온도", "-5", "35", "80", "℃"],
     "expected": "mismatch", "bad": ["upper_limit"], "reason": "기준 상한 85 vs 대상 80"},
    {"entity": "방전환경온도", "where": "s2", "row": ["방전환경온도", "-30", "30", "90", "℃"],
     "expected": "match", "reason": "-30/30/90 일치"},
    {"entity": "평가환경온도", "where": "s2", "row": ["평가환경온도", "21", "24", "28", "℃"],
     "expected": "match", "reason": "21/24/28 일치"},
    {"entity": "표준환경온도", "where": "s2-en",
     "text": "Standard ambient temperature: 21 - 29 degC (target 25 degC)",
     "expected": "match", "reason": "영문 표기지만 21/25/29 일치 — 교차언어 매칭 케이스"},
    # slide 3 — 저장 조건
    {"entity": "1개월저장온도", "where": "s3", "text": "1개월 저장: -10 ~ 80℃ (기준 35℃)",
     "expected": "match", "reason": "-10/35/80 일치"},
    {"entity": "3개월저장온도", "where": "s3", "text": "3개월 저장: -10 ~ 60℃ (기준 35℃)",
     "expected": "mismatch", "bad": ["upper_limit"], "reason": "기준 상한 70 vs 대상 60"},
    {"entity": "1년저장온도", "where": "s3", "text": "1년 저장: -10 ~ 55℃ (기준 35℃)",
     "expected": "match", "reason": "-10/35/55 일치"},
    # slide 4 — 본문 + 스피커노트로 정보가 나뉜 케이스
    {"entity": "정격충전전압", "where": "s4", "text": "충전 전압 4.55V",
     "expected": "match", "reason": "본문 값 일치(측정조건은 스피커노트에 분리 기재)"},
    {"entity": "공칭전압", "where": "s4-note", "text": "공칭전압 3.89V, 0.1C 조건 기준",
     "expected": "match",
     "reason": "값이 **스피커노트에만** 존재 — 노트를 못 읽으면 missing 오판이 나는 케이스"},
    # 대상에 없는 항목
    {"entity": "공칭용량", "where": "absent", "expected": "missing",
     "reason": "발표자료에 용량 언급 없음"},
    {"entity": "최대충전전류", "where": "absent", "expected": "missing",
     "reason": "발표자료에 전류 언급 없음"},
    {"entity": "deltaOCV", "where": "absent", "expected": "missing",
     "reason": "발표자료에 언급 없음"},
]


# --------------------------------------------------------------------------- #
# 문서 생성
# --------------------------------------------------------------------------- #
def build_word(path: str) -> None:
    """문단 + 표가 섞인 Word 대상 문서. raw/word_raw.py 의 두 블록 경로를 모두 태운다."""
    from docx import Document

    doc = Document()
    doc.add_heading("배터리 셀 규격서 (요약)", level=1)
    doc.add_paragraph("본 문서는 자표준문서의 규격 항목을 서술형으로 정리한 것이다.")

    doc.add_heading("1. 전기적 규격", level=2)
    for case in WORD_CASES:
        if case["where"] == "para":
            doc.add_paragraph(case["text"])

    doc.add_heading("2. 주요 사양 표", level=2)
    table_cases = [c for c in WORD_CASES if c["where"] == "table"]
    table = doc.add_table(rows=1 + len(table_cases), cols=3)
    table.style = "Table Grid"
    for i, head in enumerate(["항목", "규격", "단위"]):
        table.cell(0, i).text = head
    for r, case in enumerate(table_cases, start=1):
        for c, val in enumerate(case["row"]):
            table.cell(r, c).text = val

    doc.add_paragraph("※ 본 규격서에 명시되지 않은 항목은 별도 협의한다.")
    doc.save(path)


def build_ppt(path: str) -> None:
    """텍스트박스 + 표 + 스피커노트를 쓰는 PPT 대상 문서.

    슬라이드 4 는 본문과 스피커노트에 정보를 **나눠** 담아, F3 이 슬라이드 단위로
    본문+노트를 하나의 fact 로 병합하는지 검증한다(계획 §1.2 의 핵심 예시).
    """
    from pptx import Presentation
    from pptx.util import Cm, Pt

    prs = Presentation()
    blank = prs.slide_layouts[6]  # 빈 레이아웃 — 모든 도형을 직접 배치

    def textbox(slide, text: str, top_cm: float, size: int = 18):
        box = slide.shapes.add_textbox(Cm(1.5), Cm(top_cm), Cm(22), Cm(2))
        frame = box.text_frame
        frame.word_wrap = True
        frame.text = text
        frame.paragraphs[0].runs[0].font.size = Pt(size)
        return box

    def by(where: str) -> list[dict]:
        return [c for c in PPT_CASES if c["where"] == where]

    # slide 1 — 개요
    s1 = prs.slides.add_slide(blank)
    textbox(s1, "배터리 셀 규격 리뷰", 2.0, size=32)
    for i, case in enumerate(by("s1")):
        textbox(s1, case["text"], 6.0 + i * 2.0)

    # slide 2 — 환경 조건 표 + 영문 표기 텍스트
    s2 = prs.slides.add_slide(blank)
    textbox(s2, "환경 조건", 1.0, size=28)
    rows = by("s2")
    table = s2.shapes.add_table(
        len(rows) + 1, 5, Cm(1.5), Cm(3.5), Cm(22), Cm(1.2 * (len(rows) + 1))
    ).table
    for i, head in enumerate(["항목", "하한", "중심", "상한", "단위"]):
        table.cell(0, i).text = head
    for r, case in enumerate(rows, start=1):
        for c, val in enumerate(case["row"]):
            table.cell(r, c).text = val
    for i, case in enumerate(by("s2-en")):
        textbox(s2, case["text"], 12.0 + i * 1.5, size=14)
    s2.notes_slide.notes_text_frame.text = "환경 조건은 챔버 기준, 셀 표면 온도 아님."

    # slide 3 — 저장 조건
    s3 = prs.slides.add_slide(blank)
    textbox(s3, "저장 조건", 1.0, size=28)
    for i, case in enumerate(by("s3")):
        textbox(s3, case["text"], 3.5 + i * 2.0)
    s3.notes_slide.notes_text_frame.text = "저장 조건은 SOC 50% 기준."

    # slide 4 — 본문/노트 분리
    s4 = prs.slides.add_slide(blank)
    textbox(s4, "충전 규격", 1.0, size=28)
    for i, case in enumerate(by("s4")):
        textbox(s4, case["text"], 3.5 + i * 2.0)
    notes = [c["text"] for c in by("s4-note")]
    s4.notes_slide.notes_text_frame.text = " / ".join(notes + ["측정 조건: 0.1C, 4.55V"])

    prs.save(path)


# --------------------------------------------------------------------------- #
# 골든셋
# --------------------------------------------------------------------------- #
def build_golden() -> list[dict]:
    """CASES → 골든 레코드(entity × 대상문서 단위). F5 비교 결과 스키마와 1:1 대응."""
    out: list[dict] = []
    for doc_name, cases, tag in (
        (WORD_NAME, WORD_CASES, "word"),
        (PPT_NAME, PPT_CASES, "ppt"),
    ):
        for i, case in enumerate(cases, start=1):
            entity = case["entity"]
            ref = REF[entity]
            out.append({
                "id": f"g-{tag}-{i:02d}",
                "entity_name": entity,
                "target_doc": doc_name,
                "reference": {
                    "doc": "자표준문서.xlsx",
                    "row": ref["row"],
                    "cell_range": ref["cells"],
                    "attributes": ref["attrs"],
                },
                "target_text": case.get("text") or (
                    " | ".join(case["row"]) if case.get("row") else ""
                ),
                "expected": case["expected"],
                "mismatch_attributes": case.get("bad", []),
                "reason": case["reason"],
            })
    return out


def main(argv: list[str] | None = None) -> int:
    p = argparse.ArgumentParser(description="자표준문서용 합성 대상 문서 + 골든셋 생성")
    p.add_argument("--out-dir", default=".", help="저장 루트(기본: 저장소 루트)")
    args = p.parse_args(argv)

    root = os.path.abspath(args.out_dir)
    samples = os.path.join(root, "samples")
    golden_dir = os.path.join(root, "golden")
    os.makedirs(samples, exist_ok=True)
    os.makedirs(golden_dir, exist_ok=True)

    word_path = os.path.join(samples, WORD_NAME)
    ppt_path = os.path.join(samples, PPT_NAME)
    golden_path = os.path.join(golden_dir, "자표준_골든셋.jsonl")

    build_word(word_path)
    build_ppt(ppt_path)
    records = build_golden()
    with open(golden_path, "w", encoding="utf-8", newline="\n") as f:
        for rec in records:
            f.write(json.dumps(rec, ensure_ascii=False) + "\n")

    counts: dict[str, int] = {}
    for rec in records:
        counts[rec["expected"]] = counts.get(rec["expected"], 0) + 1
    print(f"생성: {word_path}")
    print(f"생성: {ppt_path}")
    print(f"생성: {golden_path} ({len(records)}항목 {counts})")
    return 0


if __name__ == "__main__":
    sys.exit(main())
